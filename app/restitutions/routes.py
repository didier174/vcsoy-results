"""
Module « Restitution ».

Même principe que « Rapport d'étude » (voir app/reports/routes.py) : liste
des présentations de restitution existantes (nom + date de création),
chargement d'un modèle de présentation (.pptx, stocké en base comme les
rapports), création d'une restitution pour un participant à partir d'un
modèle (balises {{ ... }} remplacées par ses données, même moteur que le
rapport d'étude — voir reports/generator.py et reports/report_data.py),
chargement direct d'une restitution déjà prête depuis le disque local,
téléchargement et suppression de restitutions.

Le futur enrichissement (identification par canal des 3 tests les moins
bien notés, récupération de leurs records, caviardage et insertion dans la
présentation) sera ajouté une fois le modèle de présentation finalisé.
"""

import io
import mimetypes
import re

from flask import Blueprint, current_app, render_template, request, redirect, url_for, flash, send_file
from flask_login import login_required, current_user
from pptx import Presentation

from app.access_control import admin_required
from app.extensions import db
from app.models import (
    Category, Restitution, RestitutionTemplate, Participant, TestResult, TestRecord, RedactedRecord,
    RestitutionTestSelection, ActionLog,
)
from app.editions import get_current_edition_id, get_edition
from app.menu import MENU_ITEMS
from app.reports.generator import substitute_tags
from app.reports.report_cache import get_fresh_edition_cache
from app.reports.report_data import build_participant_placeholders
from app.restitutions.debrief_visuals import apply_debrief_visuals, CRITERION_SHORT_NAMES
from app.restitutions.redaction import redact_pdf
from app.results.presentation import CHANNEL_LABELS, CHANNEL_ORDER, build_test_view
from app.results.scoring import compute_test_score, is_test_completed, CRITERIA_BY_CHANNEL
from app.results.validation import CHANNEL_FIELD_BY_KEY, CHANNELS

# {canal: (numéro mini, numéro maxi)} de la plage fixe encodée dans le
# numéro de test (voir CHANNELS ci-dessus, autorité déjà utilisée pour
# valider les numéros de test à l'import des résultats) — réutilisée ici
# pour la sélection directe d'un test par son numéro complet.
CHANNEL_ID_RANGE = {v["key"]: v["range"] for v in CHANNELS.values()}

restitutions_bp = Blueprint("restitutions", __name__, url_prefix="/restitutions")

ACTIVE_ITEM = "Restitution"
ACTIVE_ITEM_SELECT_TESTS = "Sélection des tests pour restitution"
ACTIVE_ITEM_REDACT_RECORDS = "Caviarder des records"

MAX_TESTS_REQUEST = 5

# Accumule les tests trouvés au fil de PLUSIEURS recherches (canaux/critères
# différents), pour qu'ils restent affichés tant qu'ils ne sont pas
# explicitement supprimés ou la sélection vidée — une nouvelle recherche
# s'AJOUTE à la liste au lieu de la remplacer (voir demande explicite de
# l'utilisateur : "tu ne conserves pas affichée la liste des tests
# sélectionnés"). Persistée en base (RestitutionTestSelection), pas en
# session Flask, pour ne plus disparaître d'une connexion à l'autre (voir
# demande explicite : "la sélection disparaît d'une session à l'autre de
# l'outil") — propre à chaque utilisateur et à chaque édition.


def _get_selection_entries(edition_id):
    """{test_id (str): {"note": float, "codes": [int]|None|"manuel"}} —
    codes=None signifie "tous les critères du canal"."""
    rows = RestitutionTestSelection.query.filter_by(user_id=current_user.id, edition_id=edition_id).all()
    return {str(r.test_result_id): {"note": r.note_20, "codes": r.codes} for r in rows}


def _set_selection_entries(edition_id, entries_by_id):
    """Remplace intégralement la sélection accumulée par `entries_by_id`
    (même contrat que l'ancienne version en session : les appelants
    relisent puis réécrivent le dict complet)."""
    RestitutionTestSelection.query.filter_by(user_id=current_user.id, edition_id=edition_id).delete(
        synchronize_session=False
    )
    for id_str, entry in entries_by_id.items():
        db.session.add(RestitutionTestSelection(
            user_id=current_user.id, edition_id=edition_id, test_result_id=int(id_str),
            codes=entry.get("codes"), note_20=entry.get("note"),
        ))
    db.session.commit()


def _format_codes_display(codes):
    """"Sélection manuelle" pour un test choisi directement par son numéro,
    "Tous" si aucun critère spécifique (tous les critères du canal),
    sinon "2, 3 et 6" (voir demande explicite de l'utilisateur)."""
    if codes == "manuel":
        return "Sélection manuelle"
    if not codes:
        return "Tous"
    codes = sorted(codes)
    if len(codes) == 1:
        return str(codes[0])
    return ", ".join(str(c) for c in codes[:-1]) + " et " + str(codes[-1])

RESTITUTION_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

FILENAME_UNSAFE_RE = re.compile(r'[\\/:*?"<>|]')


def _sanitize_restitution_filename(raw):
    """
    Nettoie le nom de fichier saisi par l'utilisateur : retire une
    éventuelle extension .pptx déjà tapée (l'extension est toujours
    imposée par le serveur) et remplace les caractères interdits dans un
    nom de fichier, sans toucher aux accents/espaces.
    """
    name = re.sub(r"\.pptx$", "", raw.strip(), flags=re.IGNORECASE)
    name = FILENAME_UNSAFE_RE.sub("_", name).strip()
    return name or "Restitution"


def _log(action, details=""):
    entry = ActionLog(
        user_id=current_user.id, user_email=current_user.email,
        edition_id=get_current_edition_id(), action=action, details=details,
    )
    db.session.add(entry)
    db.session.commit()


@restitutions_bp.route("/", methods=["GET"])
@login_required
def list_restitutions():
    edition_id = get_current_edition_id()
    # defer(file_data) : voir reports.list_reports, même raison (éviter de
    # charger en mémoire le contenu binaire complet de toutes les
    # présentations déjà générées à chaque affichage de cette liste).
    restitutions = (
        Restitution.query.options(db.defer(Restitution.file_data))
        .filter_by(edition_id=edition_id).order_by(Restitution.created_at.desc()).all()
    )
    templates = (
        RestitutionTemplate.query.options(db.defer(RestitutionTemplate.file_data))
        .filter_by(edition_id=edition_id).order_by(RestitutionTemplate.uploaded_at.desc()).all()
    )
    participants = Participant.query.filter_by(edition_id=edition_id).order_by(Participant.participant_name).all()
    edition = get_edition(edition_id)
    return render_template(
        "restitutions/list.html", edition=edition, restitutions=restitutions, templates=templates,
        participants=participants, active_item=ACTIVE_ITEM, menu_items=MENU_ITEMS,
    )


def _participant_channel_counts(edition_id, participant_id):
    """{canal: nb de tests} pour ce participant, tous canaux confondus."""
    rows = (
        db.session.query(TestResult.channel, db.func.count(TestResult.id))
        .filter_by(edition_id=edition_id, participant_id=participant_id)
        .group_by(TestResult.channel)
        .all()
    )
    return dict(rows)


def _search_worst_tests(edition_id, args):
    """Cherche, pour un participant/canal donnés, les `nb_tests` tests
    complets (QS completed) les moins bien notés — sur tous les critères
    du canal, ou seulement sur la sélection de critères demandée.

    Retourne (entries_by_id, error) : `entries_by_id` est un dict
    {test_id: {"note": note_20, "codes": codes}} des tests trouvés
    (destiné à être fusionné dans la sélection accumulée en session, pas
    affiché directement — `codes` est None si "tous les critères" a été
    demandé) ; `error` est un message à afficher si la sélection est
    invalide ou si aucun test ne correspond."""
    participant_id = args.get("participant_id", "").strip()
    channel = args.get("channel", "").strip()
    nb_tests_raw = args.get("nb_tests", "").strip()
    criteria_mode = args.get("criteria_mode", "tous").strip()
    codes_raw = [c for c in args.getlist("codes") if c.isdigit()]

    if not participant_id.isdigit() or channel not in CRITERIA_BY_CHANNEL:
        return {}, "Sélection invalide : choisissez un participant et un canal."

    try:
        nb_tests = int(nb_tests_raw)
    except ValueError:
        nb_tests = 0
    if not (1 <= nb_tests <= MAX_TESTS_REQUEST):
        return {}, f"Le nombre de tests souhaité doit être entre 1 et {MAX_TESTS_REQUEST}."

    participant = Participant.query.filter_by(id=int(participant_id), edition_id=edition_id).first()
    if not participant:
        return {}, "Participant introuvable pour cette édition."

    codes = None
    if criteria_mode == "selection":
        codes = [int(c) for c in codes_raw]
        if not codes:
            return {}, "Merci de choisir au moins un critère, ou « Tous les critères »."

    tests = TestResult.query.filter_by(edition_id=edition_id, participant_id=participant.id, channel=channel).all()
    scored = []
    for t in tests:
        raw_data = t.raw_data or {}
        if not is_test_completed(channel, raw_data):
            continue
        score = compute_test_score(channel, raw_data, codes=codes)
        if score is None:
            continue
        scored.append((score["note_20"], t))

    if not scored:
        return {}, (
            "Aucun test complet (QS completed) ne correspond à cette sélection de critères pour ce "
            "participant et ce canal."
        )

    scored.sort(key=lambda pair: pair[0])
    return {t.id: {"note": note_20, "codes": codes} for note_20, t in scored[:nb_tests]}, None


def _direct_test_lookup(edition_id, args):
    """Sélectionne UN test précis par son numéro complet (catégorie +
    participant + canal + numéro dans la plage fixe du canal — voir
    CHANNEL_ID_RANGE), reconstitué à partir des menus déroulants plutôt que
    saisi librement (voir demande explicite : impossible de composer un
    numéro hors format). Aucun filtre de complétude ni de critère —
    l'utilisateur choisit ce test précisément, pas "le pire selon...".

    Retourne (entries_by_id, error) — même forme que _search_worst_tests."""
    participant_id = args.get("direct_participant_id", "").strip()
    channel = args.get("direct_channel", "").strip()
    number_raw = args.get("direct_number", "").strip()

    if not participant_id.isdigit() or channel not in CHANNEL_ID_RANGE:
        return {}, "Sélection invalide : choisissez un participant et un canal."

    participant = Participant.query.filter_by(id=int(participant_id), edition_id=edition_id).first()
    if not participant:
        return {}, "Participant introuvable pour cette édition."

    lo, hi = CHANNEL_ID_RANGE[channel]
    if not number_raw.isdigit() or not (lo <= int(number_raw) <= hi):
        return {}, f"Numéro de test invalide pour ce canal (attendu entre {lo} et {hi})."

    category_code = (participant.category.code if participant.category else "").strip()
    participant_code = (participant.code or "").strip()
    if not re.fullmatch(r"\d{2}", category_code) or not re.fullmatch(r"\d{2}", participant_code):
        return {}, "Code catégorie ou participant invalide pour ce participant."

    test_id_str = f"{category_code}{participant_code}{int(number_raw):04d}"
    t = TestResult.query.filter_by(edition_id=edition_id, test_id=test_id_str, channel=channel).first()
    if not t:
        return {}, f"Aucun test {test_id_str} chargé pour ce participant et ce canal."

    score = compute_test_score(channel, t.raw_data or {})
    return {t.id: {"note": score["note_20"] if score else None, "codes": "manuel"}}, None


def _load_selection_results(edition_id):
    """Reconstruit la liste affichable (voir build_test_view) de la
    sélection accumulée en session — partagée entre "Sélection des tests
    pour restitution" et "Caviarder des records", les deux pages
    travaillant sur la MÊME sélection. Nettoie au passage les entrées dont
    le test a été supprimé entretemps (par un autre onglet, etc.)."""
    entries_by_id = _get_selection_entries(edition_id)
    if not entries_by_id:
        return []

    selected_ids = [int(i) for i in entries_by_id]
    tests_by_id = {
        t.id: t for t in TestResult.query.filter(
            TestResult.edition_id == edition_id, TestResult.id.in_(selected_ids)
        ).all()
    }
    results = []
    still_valid = {}
    for id_str, entry in entries_by_id.items():
        t = tests_by_id.get(int(id_str))
        if t is None:
            continue  # supprimé entretemps
        view = build_test_view(t)
        view["note_20"] = entry["note"]
        view["codes_display"] = _format_codes_display(entry.get("codes"))
        results.append(view)
        still_valid[id_str] = entry
    if still_valid.keys() != entries_by_id.keys():
        _set_selection_entries(edition_id, still_valid)
    # un test choisi directement par son numéro peut n'avoir aucun critère
    # valide renseigné (note_20 = None) — placé en dernier plutôt que de
    # faire échouer le tri (None n'est pas comparable à un nombre).
    results.sort(key=lambda v: (v["note_20"] is None, v["note_20"]))
    return results


@restitutions_bp.route("/selection-tests", methods=["GET"])
@login_required
def select_tests():
    """Choix, par canal, des tests les moins bien notés à insérer (records
    caviardés) dans la restitution. La liste des participants proposés, et
    pour chacun les canaux proposés, sont déjà filtrés aux seuls canaux
    déclarés actifs POUR CE PARTICIPANT et ayant au moins un test chargé
    dans l'édition — ce qui vérifie par construction, avant même d'ouvrir
    le pop-up, qu'une recherche pourra aboutir (voir demande explicite de
    l'utilisateur).

    Chaque recherche (voir _search_worst_tests) AJOUTE ses résultats à la
    sélection déjà accumulée en session au lieu de la remplacer, pour
    pouvoir composer une restitution à partir de plusieurs canaux sans
    perdre les tests déjà trouvés (voir demande explicite de l'utilisateur :
    la liste ne restait pas affichée d'une recherche à l'autre)."""
    edition_id = get_current_edition_id()
    edition = get_edition(edition_id)

    has_any_test = db.session.query(TestResult.id).filter_by(edition_id=edition_id).first() is not None

    participants = Participant.query.filter_by(edition_id=edition_id).order_by(Participant.participant_name).all()
    participant_channels = {}
    participant_categories = {}
    for p in participants:
        counts = _participant_channel_counts(edition_id, p.id)
        valid = [c for c in CHANNEL_ORDER if getattr(p, CHANNEL_FIELD_BY_KEY[c], False) and counts.get(c, 0) > 0]
        if valid:
            participant_channels[p.id] = valid
            participant_categories[p.id] = p.category_id
    participants_with_tests = [p for p in participants if p.id in participant_channels]
    categories = (
        Category.query.filter_by(edition_id=edition_id)
        .filter(Category.id.in_(set(participant_categories.values())))
        .order_by(Category.code).all()
    ) if participant_categories else []

    criteria_by_channel = {
        channel: [{"code": code, "name": CRITERION_SHORT_NAMES[channel][code]} for code in sorted(codes)]
        for channel, codes in CRITERIA_BY_CHANNEL.items()
    }

    mode = request.args.get("mode", "search")
    if request.args.get("participant_id") or request.args.get("direct_participant_id"):
        if mode == "direct":
            new_entries, error = _direct_test_lookup(edition_id, request.args)
        else:
            new_entries, error = _search_worst_tests(edition_id, request.args)
        if error:
            flash(error, "error")
        if new_entries:
            entries_by_id = _get_selection_entries(edition_id)
            entries_by_id.update({str(tid): entry for tid, entry in new_entries.items()})
            _set_selection_entries(edition_id, entries_by_id)
        # redirection (PRG) : évite de relancer la même recherche à
        # chaque rafraîchissement de la page, et affiche l'URL "propre".
        return redirect(url_for("restitutions.select_tests"))

    results = _load_selection_results(edition_id) or None

    return render_template(
        "restitutions/select_tests.html", edition=edition,
        active_item=ACTIVE_ITEM_SELECT_TESTS, menu_items=MENU_ITEMS,
        has_any_test=has_any_test, participants=participants_with_tests,
        participant_channels=participant_channels, participant_categories=participant_categories,
        categories=categories, channel_labels=CHANNEL_LABELS, channel_id_ranges=CHANNEL_ID_RANGE,
        criteria_by_channel=criteria_by_channel, max_tests=MAX_TESTS_REQUEST,
        results=results,
    )


@restitutions_bp.route("/selection-tests/clear", methods=["POST"])
@login_required
def clear_selected_tests():
    """Vide la sélection accumulée (sans toucher aux tests en base — voir
    delete_selected_tests pour la suppression définitive)."""
    edition_id = get_current_edition_id()
    _set_selection_entries(edition_id, {})
    flash("Sélection réinitialisée.", "success")
    return redirect(url_for("restitutions.select_tests"))


@restitutions_bp.route("/selection-tests/delete", methods=["POST"])
@login_required
@admin_required
def delete_selected_tests():
    """Supprime les tests sélectionnés dans le tableau de résultats (et
    leur record associé, s'il existe) — même principe que
    records.delete_records : le record doit être supprimé avant le test
    (pas de suppression en cascade configurée sur la relation)."""
    edition_id = get_current_edition_id()
    test_ids = [int(i) for i in request.form.getlist("test_result_ids") if i.isdigit()]
    if not test_ids:
        flash("Merci de choisir au moins un test à supprimer.", "error")
    else:
        to_delete = TestResult.query.filter(
            TestResult.edition_id == edition_id, TestResult.id.in_(test_ids)
        ).all()
        deleted = len(to_delete)
        ids = [t.id for t in to_delete]
        TestRecord.query.filter(TestRecord.test_result_id.in_(ids)).delete(synchronize_session=False)
        TestResult.query.filter(TestResult.edition_id == edition_id, TestResult.id.in_(ids)).delete(
            synchronize_session=False
        )
        db.session.commit()

        entries_by_id = _get_selection_entries(edition_id)
        remaining = {k: v for k, v in entries_by_id.items() if int(k) not in ids}
        if remaining.keys() != entries_by_id.keys():
            _set_selection_entries(edition_id, remaining)

        _log("Suppression de test(s) (sélection restitution)", details=f"{deleted} supprimé(s) (édition {edition_id})")
        flash(f"{deleted} test(s) supprimé(s).", "success")

    return redirect(url_for("restitutions.select_tests"))


def _annotate_caviardage_eligibility(results):
    """Ajoute à chaque vue de test (voir build_test_view) : is_pdf_record
    (éligible au caviardage — record présent et pas un fichier audio, seul
    canal — Téléphone — dont le record n'est pas un PDF) et
    redacted_record_id (copie déjà caviardée, s'il y en a une)."""
    test_ids = [r["id"] for r in results]
    redacted_by_test = {
        rr.test_result_id: rr.id
        for rr in RedactedRecord.query.filter(RedactedRecord.test_result_id.in_(test_ids)).all()
    } if test_ids else {}
    for r in results:
        r["is_pdf_record"] = bool(r["record_id"]) and not r["record_is_audio"]
        r["redacted_record_id"] = redacted_by_test.get(r["id"])
    return results


@restitutions_bp.route("/caviardage", methods=["GET"])
@login_required
def redact_records():
    """Caviardage des records sélectionnés avant insertion dans la
    restitution : reprend la MÊME sélection accumulée que « Sélection des
    tests pour restitution » (voir _load_selection_results) — sans le
    bouton Détail test, avec les boutons Caviarder/Supprimer agissant
    uniquement sur les tests cochés (voir demande explicite)."""
    edition_id = get_current_edition_id()
    edition = get_edition(edition_id)
    results = _load_selection_results(edition_id)
    _annotate_caviardage_eligibility(results)
    return render_template(
        "restitutions/redact_records.html", edition=edition,
        active_item=ACTIVE_ITEM_REDACT_RECORDS, menu_items=MENU_ITEMS,
        results=results,
    )


@restitutions_bp.route("/caviardage/redact", methods=["POST"])
@login_required
def redact_selected_tests():
    """Caviarde le record de chaque test coché : duplique son contenu (le
    TestRecord d'origine n'est jamais lu en écriture) et retire du double
    le nom du participant (connu), les adresses e-mail (masquées, voir
    redaction.py), et — si renseignés dans le formulaire, laissés vides
    sinon (voir demande explicite) — le nom du testeur et celui du
    conseiller client, aucun des deux n'existant en base sous forme
    structurée."""
    edition_id = get_current_edition_id()
    test_ids = [int(i) for i in request.form.getlist("test_result_ids") if i.isdigit()]
    if not test_ids:
        flash("Merci de choisir au moins un test à caviarder.", "error")
        return redirect(url_for("restitutions.redact_records"))

    tests = TestResult.query.filter(TestResult.edition_id == edition_id, TestResult.id.in_(test_ids)).all()

    redacted, skipped_no_record, skipped_not_pdf = [], [], []
    tester_not_found, advisor_not_found = [], []
    for t in tests:
        record = t.record
        if record is None:
            skipped_no_record.append(t.test_id)
            continue
        if record.is_audio:
            skipped_not_pdf.append(t.test_id)
            continue

        participant_name = t.participant.participant_name if t.participant else None
        tester_name = request.form.get(f"tester_name_{t.id}", "").strip() or None
        advisor_name = request.form.get(f"advisor_name_{t.id}", "").strip() or None

        try:
            redacted_bytes, not_found = redact_pdf(
                record.file_data, participant_name=participant_name,
                tester_name=tester_name, advisor_name=advisor_name,
            )
        except Exception:
            current_app.logger.exception("Échec du caviardage du record du test %s", t.test_id)
            skipped_not_pdf.append(t.test_id)
            continue
        if not_found.get("testeur"):
            tester_not_found.append(t.test_id)
        if not_found.get("conseiller"):
            advisor_not_found.append(t.test_id)

        filename = re.sub(
            r"\.pdf$", "_caviarde.pdf", record.filename or f"{t.test_id}-record.pdf", flags=re.IGNORECASE
        )
        existing = RedactedRecord.query.filter_by(test_result_id=t.id).first()
        if existing:
            existing.filename = filename
            existing.content_type = "application/pdf"
            existing.file_data = redacted_bytes
            existing.file_size = len(redacted_bytes)
            existing.created_by_id = current_user.id
            record_obj = existing
        else:
            record_obj = RedactedRecord(
                test_result_id=t.id, filename=filename, content_type="application/pdf",
                file_data=redacted_bytes, file_size=len(redacted_bytes), created_by_id=current_user.id,
            )
            db.session.add(record_obj)
        redacted.append(t.test_id)

        # Enregistré tout de suite (pas un seul commit à la fin de la
        # boucle) puis retiré de la session : sur un gros lot (édition à
        # venir, ~250-300 records à caviarder d'un coup), garder tous les
        # octets déjà caviardés en mémoire jusqu'à la fin de la requête
        # gonflerait inutilement sa consommation mémoire, en plus du pic
        # transitoire du caviardage du record en cours.
        db.session.commit()
        db.session.expunge(record_obj)

    if redacted:
        _log(
            "Caviardage de record(s)",
            details=f"{len(redacted)} test(s) (édition {edition_id}) : {', '.join(redacted)}",
        )
        flash(f"{len(redacted)} record(s) caviardé(s) : {', '.join(redacted)}.", "success")
    if skipped_no_record:
        flash("Impossible de caviarder (aucun record) : " + ", ".join(skipped_no_record) + ".", "error")
    if skipped_not_pdf:
        flash(
            "Impossible de caviarder (record audio — seuls les PDF sont pris en charge pour l'instant) : "
            + ", ".join(skipped_not_pdf) + ".",
            "error",
        )
    if tester_not_found:
        flash(
            "Nom du testeur introuvable dans le record (vérifiez l'orthographe, ou le PDF ne contient peut-être "
            "pas de texte sélectionnable) — non caviardé pour : " + ", ".join(tester_not_found) + ".",
            "error",
        )
    if advisor_not_found:
        flash(
            "Nom du conseiller client introuvable dans le record (vérifiez l'orthographe, ou le PDF ne contient "
            "peut-être pas de texte sélectionnable) — non caviardé pour : " + ", ".join(advisor_not_found) + ".",
            "error",
        )

    return redirect(url_for("restitutions.redact_records"))


@restitutions_bp.route("/caviardage/delete", methods=["POST"])
@login_required
def delete_caviardage_tests():
    """Supprime la copie caviardée des tests cochés — le test lui-même
    RESTE dans la sélection (partagée avec « Sélection des tests pour
    restitution »), réaffiché comme avant la demande de caviardage (voir
    demande explicite de l'utilisateur : ce bouton ne doit pas retirer le
    test de la sélection, seulement sa copie caviardée). Ne touche jamais
    au TestResult ni à son TestRecord d'origine."""
    edition_id = get_current_edition_id()
    test_ids = [int(i) for i in request.form.getlist("test_result_ids") if i.isdigit()]
    if not test_ids:
        flash("Merci de choisir au moins un test dont retirer la copie caviardée.", "error")
        return redirect(url_for("restitutions.redact_records"))

    deleted_redacted = RedactedRecord.query.filter(RedactedRecord.test_result_id.in_(test_ids)).delete(
        synchronize_session=False
    )
    db.session.commit()

    _log(
        "Suppression de copie(s) caviardée(s)",
        details=f"{deleted_redacted} copie(s) supprimée(s) sur {len(test_ids)} test(s) sélectionné(s) (édition {edition_id})",
    )
    if deleted_redacted == 0:
        flash("Aucune copie caviardée à supprimer pour le(s) test(s) coché(s).", "error")
    else:
        flash(
            f"{deleted_redacted} copie(s) caviardée(s) supprimée(s) — le(s) test(s) reste(nt) dans la sélection.",
            "success",
        )
    return redirect(url_for("restitutions.redact_records"))


@restitutions_bp.route("/redacted/<int:redacted_id>/download", methods=["GET"])
@login_required
def download_redacted_record(redacted_id):
    edition_id = get_current_edition_id()
    record = (
        RedactedRecord.query.join(TestResult)
        .filter(RedactedRecord.id == redacted_id, TestResult.edition_id == edition_id)
        .first()
    )
    if not record:
        return "Record caviardé introuvable pour cette édition.", 404

    return send_file(
        io.BytesIO(record.file_data), mimetype=record.content_type or "application/pdf",
        as_attachment=False, download_name=record.filename,
    )


@restitutions_bp.route("/templates/upload", methods=["POST"])
@login_required
def upload_template():
    edition_id = get_current_edition_id()
    file = request.files.get("template_file")
    if not file or not file.filename:
        flash("Merci de choisir un fichier avant de cliquer sur « Charger ».", "error")
        return redirect(url_for("restitutions.list_restitutions"))

    filename = file.filename
    content = file.read()
    content_type = file.content_type or mimetypes.guess_type(filename)[0] or "application/octet-stream"

    db.session.add(RestitutionTemplate(
        edition_id=edition_id, filename=filename, content_type=content_type,
        file_data=content, file_size=len(content), uploaded_by_id=current_user.id,
    ))
    db.session.commit()

    _log("Chargement d'un modèle de restitution", details=f"{filename} (édition {edition_id})")
    flash(f"Modèle « {filename} » chargé avec succès.", "success")
    return redirect(url_for("restitutions.list_restitutions"))


@restitutions_bp.route("/templates/delete", methods=["POST"])
@login_required
@admin_required
def delete_templates():
    edition_id = get_current_edition_id()
    selected_ids = [int(i) for i in request.form.getlist("template_ids") if i.isdigit()]
    if not selected_ids:
        flash("Merci de choisir au moins un modèle à supprimer.", "error")
        return redirect(url_for("restitutions.list_restitutions"))

    to_delete = RestitutionTemplate.query.options(db.defer(RestitutionTemplate.file_data)).filter(
        RestitutionTemplate.edition_id == edition_id, RestitutionTemplate.id.in_(selected_ids)
    ).all()
    deleted = len(to_delete)
    names = ", ".join(t.filename for t in to_delete)

    ids = [t.id for t in to_delete]
    Restitution.query.filter(Restitution.restitution_template_id.in_(ids)).update(
        {"restitution_template_id": None}, synchronize_session=False
    )
    for template in to_delete:
        db.session.delete(template)
    db.session.commit()

    _log("Suppression de modèle(s) de restitution", details=f"{deleted} supprimé(s) (édition {edition_id}) : {names}")
    flash(f"{deleted} modèle(s) supprimé(s).", "success")
    return redirect(url_for("restitutions.list_restitutions"))


@restitutions_bp.route("/new", methods=["POST"])
@login_required
def create_restitution():
    edition_id = get_current_edition_id()
    template_id = request.form.get("template_id", "").strip()
    participant_id = request.form.get("participant_id", "").strip()
    restitution_filename = request.form.get("restitution_filename", "").strip()

    if not template_id.isdigit() or not participant_id.isdigit() or not restitution_filename:
        flash("Merci de choisir un modèle, un participant et un nom de fichier.", "error")
        return redirect(url_for("restitutions.list_restitutions"))

    template = RestitutionTemplate.query.filter_by(id=int(template_id), edition_id=edition_id).first()
    participant = Participant.query.filter_by(id=int(participant_id), edition_id=edition_id).first()
    if not template or not participant:
        flash("Modèle ou participant introuvable pour cette édition.", "error")
        return redirect(url_for("restitutions.list_restitutions"))

    # Même passerelle que le rapport d'étude (voir reports.create_report) :
    # la restitution compare elle aussi chaque participant à l'ensemble des
    # résultats de l'édition, via le même cache d'agrégats précalculé par
    # « Liste des lauréats » — pour ne pas recharger/retraiter tous les
    # tests de l'édition à chaque présentation générée.
    cache = get_fresh_edition_cache(edition_id)
    if cache is None or not cache["has_winners"]:
        flash(
            "Aucun lauréat n'a pu être déterminé pour cette édition (aucun "
            "résultat chargé, ou aucun participant n'atteint 11,5/20 et le "
            "1er rang de sa catégorie), ou les résultats ont changé depuis "
            "le dernier calcul. La restitution ne peut pas être générée : "
            "chargez les fichiers de résultats de toute l'édition, vérifiez "
            "« Liste des lauréats », puis réessayez.",
            "error",
        )
        return redirect(url_for("restitutions.list_restitutions"))

    # Seuls les tests de CE participant sont chargés ici (pas ceux de toute
    # l'édition) : le reste vient du cache ci-dessus.
    participant_tests = TestResult.query.filter_by(edition_id=edition_id, participant_id=participant.id).all()
    values = build_participant_placeholders(participant, edition_id, cache, all_tests=participant_tests)

    prs = Presentation(io.BytesIO(template.file_data))

    unknown_tags = substitute_tags(prs, values)

    if unknown_tags:
        flash(
            "Ce modèle contient des balises non reconnues : {{ "
            + " }}, {{ ".join(sorted(unknown_tags))
            + " }}. La restitution n'a pas été générée.",
            "error",
        )
        return redirect(url_for("restitutions.list_restitutions"))

    # Graphiques de comparaison/classement + répartition verte/rouge des
    # critères (voir debrief_visuals.py) : best-effort — si le modèle de
    # restitution ne correspond pas à la structure attendue (nom/position
    # de forme différente), on préfère livrer la présentation avec les
    # balises texte déjà remplies plutôt que de faire échouer toute la
    # génération.
    try:
        apply_debrief_visuals(prs, participant, cache, values)
    except Exception:
        current_app.logger.exception("Échec de la mise à jour des graphiques natifs de la restitution")

    out = io.BytesIO()
    prs.save(out)
    file_bytes = out.getvalue()

    name = _sanitize_restitution_filename(restitution_filename)
    filename = f"{name}.pptx"

    restitution = Restitution(
        edition_id=edition_id, name=name, participant_id=participant.id, restitution_template_id=template.id,
        filename=filename, content_type=RESTITUTION_CONTENT_TYPE,
        file_data=file_bytes, file_size=len(file_bytes), created_by_id=current_user.id,
    )
    db.session.add(restitution)
    db.session.commit()

    _log("Création d'une restitution", details=f"{name} (édition {edition_id})")
    flash(f"Restitution « {name} » créée avec succès.", "success")
    return redirect(url_for("restitutions.list_restitutions"))


@restitutions_bp.route("/upload", methods=["POST"])
@login_required
def upload_restitution():
    edition_id = get_current_edition_id()
    file = request.files.get("restitution_file")
    if not file or not file.filename:
        flash("Merci de choisir un fichier avant de cliquer sur « Charger ».", "error")
        return redirect(url_for("restitutions.list_restitutions"))

    if not file.filename.lower().endswith(".pptx"):
        flash("Seul le format .pptx est accepté pour une restitution.", "error")
        return redirect(url_for("restitutions.list_restitutions"))

    base_name = _sanitize_restitution_filename(file.filename)
    content = file.read()

    restitution = Restitution(
        edition_id=edition_id, name=base_name, filename=f"{base_name}.pptx",
        content_type=RESTITUTION_CONTENT_TYPE, file_data=content, file_size=len(content),
        created_by_id=current_user.id,
    )
    db.session.add(restitution)
    db.session.commit()

    _log("Chargement direct d'une restitution", details=f"{base_name} (édition {edition_id})")
    flash(f"Restitution « {base_name} » chargée avec succès.", "success")
    return redirect(url_for("restitutions.list_restitutions"))


@restitutions_bp.route("/<int:restitution_id>/download", methods=["GET"])
@login_required
def download_restitution(restitution_id):
    edition_id = get_current_edition_id()
    restitution = Restitution.query.filter_by(id=restitution_id, edition_id=edition_id).first()
    if not restitution:
        return "Restitution introuvable pour cette édition.", 404

    return send_file(
        io.BytesIO(restitution.file_data), mimetype=restitution.content_type or RESTITUTION_CONTENT_TYPE,
        as_attachment=True, download_name=restitution.filename or f"{restitution.name}.pptx",
    )


@restitutions_bp.route("/delete", methods=["POST"])
@login_required
@admin_required
def delete_restitutions():
    edition_id = get_current_edition_id()
    selected_ids = [int(i) for i in request.form.getlist("selected_ids") if i.isdigit()]
    if not selected_ids:
        flash("Merci de choisir au moins une restitution à supprimer.", "error")
        return redirect(url_for("restitutions.list_restitutions"))

    to_delete = Restitution.query.options(db.defer(Restitution.file_data)).filter(
        Restitution.edition_id == edition_id, Restitution.id.in_(selected_ids)
    ).all()
    deleted = len(to_delete)
    names = ", ".join(r.name for r in to_delete)
    for restitution in to_delete:
        db.session.delete(restitution)
    db.session.commit()

    _log("Suppression de restitution(s)", details=f"{deleted} supprimé(s) (édition {edition_id}) : {names}")
    flash(f"{deleted} restitution(s) supprimée(s).", "success")
    return redirect(url_for("restitutions.list_restitutions"))
