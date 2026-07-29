"""
Met à jour les éléments visuels spécifiques au modèle de « Restitution »
(présentation de debrief par participant), qui ne sont pas de simples
balises {{ ... }} texte :

- diapositives 5/6/9/12/15/18 : graphique de comparaison (Vous / Catégorie /
  Ensemble des participants / Lauréats), note globale et par canal.
- diapositives 5/6/9/12/15 (pas de graphique de classement pour le chat
  dans le modèle) : graphique de classement, une fenêtre de 8 participants
  centrée sur le participant courant (classement sur l'ensemble de
  l'édition, tous canaux confondus pour la diapo globale, canal par canal
  ensuite), avec sa barre mise en évidence à sa vraie position.
- diapositive 5 : texte "Classement Xème" (classement global de l'édition).
- diapositives 6/9/12/15/18 : répartition dynamique des critères en 2
  colonnes (verte à gauche / rouge à droite) selon que la note du critère,
  ramenée sur 20, est supérieure ou égale (vert) ou inférieure (rouge) à la
  note globale du canal pour ce participant.

Ces graphiques sont des graphiques 3D (bar3DChart), un type non supporté
par l'API haut niveau de python-pptx (chart.plots/series lèvent une
exception) : on manipule donc leur XML directement, comme pour les
graphiques du rapport d'étude (voir app/reports/report_visuals.py), sans
dépendre de cette API pour ces graphiques précis.
"""

import re
from copy import deepcopy

from lxml import etree
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from app.results.presentation import CHANNEL_ORDER
from app.results.scoring import CRITERIA_BY_CHANNEL, DOUBLED_CODE_BY_CHANNEL

# Diapositive (0-based) portant le graphique de comparaison + le graphique
# de classement pour chaque portée (note globale, puis chaque canal).
COMPARISON_SLIDE_BY_SCOPE = {"global": 4, "phone": 5, "mail": 8, "web": 11, "rs": 14, "chat": 17}

# Diapositive (0-based) portant la répartition verte/rouge des critères,
# une par canal (la diapo globale n'a pas de détail par critère).
CRITERIA_SLIDE_BY_CHANNEL = {"phone": 6, "mail": 9, "web": 12, "rs": 15, "chat": 18}

# Nom court par code de critère et par canal, repris tel quel des étiquettes
# déjà utilisées dans le mapping d'importance du rapport d'étude (voir
# report_visuals.MAPPING_POINT_CODE) — le critère "Impression générale" en
# est systématiquement absent (choix déjà fait dans ce mapping), complété
# ici à partir du libellé du récapitulatif des critères du rapport d'étude.
CRITERION_SHORT_NAMES = {
    "phone": {
        1: "Accueil par le conseiller", 2: "Accueil SVI", 3: "Utilisation d'un vocabulaire adapté",
        4: "Reformulation de la demande", 5: "Transfert", 6: "Écoute", 7: "Prise de congé",
        8: "Temps de décroché", 9: "Prise en charge du contact par un conseiller",
        10: "Accès au service désiré (SVI)", 11: "Réponse différée", 12: "Mise en attente",
        13: "Qualité de la réponse", 14: "Impression générale", 15: "Qualité du ton",
    },
    "mail": {
        1: "Expéditeur", 2: "Objet", 3: "Personnalisation de la réponse", 4: "Qualité de rédaction",
        5: "Résumé de la situation", 6: "Conclusion", 7: "Identification", 8: "Poursuite de la relation",
        9: "Temps de réponse", 10: "Accusé de réception", 11: "Qualité de la réponse",
        12: "Impression générale", 13: "Ton employé", 14: "Lisibilité",
    },
    "web": {
        1: "Visibilité", 2: "Fonctionnalité", 3: "Efficacité", 4: "Rétroactivité", 5: "Attractivité",
        6: "Identification", 7: "Réponse concise", 8: "Savoir faire", 9: "Savoir dire",
        10: "Poursuite de la relation", 11: "Qualité de la réponse", 12: "Impression générale",
        13: "Lisibilité",
    },
    "rs": {
        1: "Identification du participant", 2: "Coordonnées du Participant", 3: "Personnalisation",
        4: "Rédaction", 5: "Forme de la réponse", 6: "Prise de congé",
        7: "Identification de l'interlocuteur", 8: "Poursuite de la relation", 9: "Temps de la réponse",
        10: "Message relais", 11: "Qualité de la réponse", 12: "Fonctionnalité du lien",
        13: "Impression générale", 14: "Ton employé", 15: "Savoir Faire",
    },
    "chat": {
        1: "Visibilité", 2: "Temps de prise en charge", 3: "Réactivité", 4: "Accueil",
        5: "Personnalisation", 6: "Rédaction", 7: "Prise de congé", 8: "Forme de la réponse",
        9: "Fonctionnalité du lien", 10: "Qualité de la réponse", 11: "Lisibilité",
        12: "Historique de conversation", 13: "Impression générale", 14: "Ton employé",
    },
}

RANKING_WINDOW_SIZE = 8
HIGHLIGHT_COLOR = "00B2E6"
OTHER_COLOR = "7F7F7F"


def _as_float(value):
    """Convertit une valeur de `values` (float brut, ou chaîne déjà
    formatée par report_data._fmt_note/_fmt_pct, ex. "12,34" ou "45%") en
    float, ou None si absente/non calculable ("—")."""
    if value is None or value == "—":
        return None
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value).strip().rstrip("%").replace(",", ".")
    try:
        return float(text)
    except ValueError:
        return None


# --------------------------------------------------------- Lecture du XML

def _find_chart_ser_by_marker(slide, marker):
    """Retourne (ser, cat_strCache, val_numCache) du premier graphique de la
    diapo dont au moins une catégorie contient `marker` (ex. "Catégorie"
    pour le graphique de comparaison, "Participant" pour le classement)."""
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        try:
            cs = shape.chart._chartSpace
        except Exception:
            continue
        ser = cs.find(f".//{qn('c:ser')}")
        if ser is None:
            continue
        cat = ser.find(qn("c:cat"))
        if cat is None:
            continue
        cat_values = [v.text or "" for v in cat.findall(f".//{qn('c:v')}")]
        if any(marker in v for v in cat_values):
            strcache = cat.find(f".//{qn('c:strCache')}")
            val_el = ser.find(qn("c:val"))
            numcache = val_el.find(f".//{qn('c:numCache')}") if val_el is not None else None
            return ser, strcache, numcache
    return None, None, None


def _set_str_cache(strcache_el, values):
    if strcache_el is None:
        return
    for pt in list(strcache_el.findall(qn("c:pt"))):
        strcache_el.remove(pt)
    ptcount = strcache_el.find(qn("c:ptCount"))
    if ptcount is not None:
        ptcount.set("val", str(len(values)))
    for idx, value in enumerate(values):
        pt = etree.SubElement(strcache_el, qn("c:pt"))
        pt.set("idx", str(idx))
        v_el = etree.SubElement(pt, qn("c:v"))
        v_el.text = str(value)


def _set_num_cache(numcache_el, values):
    if numcache_el is None:
        return
    for pt in list(numcache_el.findall(qn("c:pt"))):
        numcache_el.remove(pt)
    ptcount = numcache_el.find(qn("c:ptCount"))
    if ptcount is not None:
        ptcount.set("val", str(len(values)))
    for idx, value in enumerate(values):
        pt = etree.SubElement(numcache_el, qn("c:pt"))
        pt.set("idx", str(idx))
        v_el = etree.SubElement(pt, qn("c:v"))
        v_el.text = repr(float(value))


def _set_dpt_colors(ser, n_points, highlight_idx):
    """Recolore les points du graphique (c:dPt) : `highlight_idx` en
    HIGHLIGHT_COLOR, tous les autres en OTHER_COLOR, et retire les points
    au-delà de `n_points` (édition avec moins de RANKING_WINDOW_SIZE
    participants pour ce canal)."""
    for dpt in list(ser.findall(qn("c:dPt"))):
        idx_el = dpt.find(qn("c:idx"))
        if idx_el is None:
            continue
        idx = int(idx_el.get("val"))
        if idx >= n_points:
            ser.remove(dpt)
            continue
        fill = dpt.find(f".//{qn('a:solidFill')}")
        if fill is None:
            continue
        for child in list(fill):
            fill.remove(child)
        color_el = etree.SubElement(fill, qn("a:srgbClr"))
        color_el.set("val", HIGHLIGHT_COLOR if idx == highlight_idx else OTHER_COLOR)


# --------------------------------------------------- Graphique de comparaison

def apply_comparison_charts(prs, values):
    """values : dict retourné par build_participant_placeholders (voir
    app/reports/report_data.py) — mêmes balises que le rapport d'étude,
    réutilisées telles quelles ici."""
    for scope, slide_idx in COMPARISON_SLIDE_BY_SCOPE.items():
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        ser, _, numcache = _find_chart_ser_by_marker(slide, "Catégorie")
        if numcache is None:
            continue

        if scope == "global":
            own = _as_float(values.get("Note consolidée"))
            categorie = _as_float(values.get("Global note categorie"))
            tous = _as_float(values.get("Global note tous"))
            laureats = _as_float(values.get("Global note laureats"))
        else:
            own = _as_float(values.get(f"Note sur 20 canal {scope}"))
            categorie = _as_float(values.get(f"Total {scope} note categorie"))
            tous = _as_float(values.get(f"{scope} note tous"))
            laureats = _as_float(values.get(f"{scope} note laureats"))

        _set_num_cache(numcache, [v if v is not None else 0.0 for v in (own, categorie, tous, laureats)])


# ------------------------------------------------------ Graphique de classement

def _ranking_window(ranking, participant_id, size=RANKING_WINDOW_SIZE):
    """Fenêtre de `size` entrées de `ranking` (déjà trié décroissant) qui
    inclut TOUJOURS le participant à sa vraie position relative — centrée
    sur lui, sauf en bord de classement (tête ou queue de fenêtre)."""
    n = len(ranking)
    pos = next((i for i, r in enumerate(ranking) if r["participant_id"] == participant_id), None)
    if pos is None or n == 0:
        return [], None
    size = min(size, n)
    half = size // 2
    start = max(0, min(pos - half, n - size))
    window = ranking[start:start + size]
    return window, pos - start


def apply_ranking_charts(prs, participant, cache):
    ranking_by_scope = cache.get("ranking", {})
    for scope, slide_idx in COMPARISON_SLIDE_BY_SCOPE.items():
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        ser, strcache, numcache = _find_chart_ser_by_marker(slide, "Participant")
        if ser is None:
            continue  # pas de graphique de classement pour ce canal (ex. chat)

        window, own_idx = _ranking_window(ranking_by_scope.get(scope, []), participant.id)
        if not window:
            continue

        labels = [
            participant.participant_name if i == own_idx else f"Participant {i + 1}"
            for i in range(len(window))
        ]
        _set_str_cache(strcache, labels)
        _set_num_cache(numcache, [w["score"] for w in window])
        _set_dpt_colors(ser, len(window), own_idx)


def _ordinal_fr(n):
    return f"{n}er" if n == 1 else f"{n}ème"


def apply_classement_text(prs, participant, cache):
    """Diapositive 5 (globale) : remplace le texte "Classement XXème" par
    le classement réel du participant sur l'ensemble des participants de
    l'édition (note consolidée)."""
    slide_idx = COMPARISON_SLIDE_BY_SCOPE["global"]
    if slide_idx >= len(prs.slides):
        return
    slide = prs.slides[slide_idx]
    ranking = cache.get("ranking", {}).get("global", [])
    pos = next((i for i, r in enumerate(ranking) if r["participant_id"] == participant.id), None)
    if pos is None:
        return

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        cnv_pr = shape._element.find(f".//{qn('p:cNvPr')}")
        if cnv_pr is not None and cnv_pr.get("hidden") == "1":
            continue  # formes de brouillon invisibles (ex. "Classement / <4"), pas le vrai texte
        if "Classement XX" not in shape.text_frame.text:
            continue
        paragraph = shape.text_frame.paragraphs[0]
        if not paragraph.runs:
            continue
        paragraph.runs[0].text = f"Classement {_ordinal_fr(pos + 1)} "
        for run in paragraph.runs[1:]:
            run.text = ""
        break


# --------------------------------------------------- Répartition des critères

CRITERION_LINE_RE = re.compile(r"^C\d+\s")


def _find_criteria_boxes(slide):
    """Les 2 encarts de critères (vert à gauche / rouge à droite) sont les
    2 formes automatiques non vides les plus BASSES de la diapo (les 2
    encarts statistiques "XX%" sont toujours positionnés au-dessus) —
    identification par position plutôt que par contenu texte actuel,
    puisque l'encart rose peut être vide dans le modèle d'origine (ex.
    Internet : simple "-")."""
    candidates = [
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and getattr(s, "has_text_frame", False)
        and s.text_frame.text.strip()
    ]
    candidates.sort(key=lambda s: s.top, reverse=True)
    boxes = candidates[:2]
    if len(boxes) != 2:
        return None, None
    boxes.sort(key=lambda s: s.left)
    return boxes[0], boxes[1]  # (vert/gauche, rouge/droite)


def _rebuild_criteria_list(shape, items):
    """items : liste de (code, nom_court). Reconstruit les paragraphes de
    la forme en clonant la mise en forme (2 runs : code en gras + nom,
    couleur déjà réglée dans le modèle) de son PREMIER paragraphe existant
    — le nombre de critères affichés varie par participant, donc le nombre
    de paragraphes doit être reconstruit à chaque génération."""
    tf = shape.text_frame
    paragraphs = tf.paragraphs
    if not paragraphs:
        return
    template_el = deepcopy(paragraphs[0]._p)
    txbody = tf._txBody
    for p in list(paragraphs):
        txbody.remove(p._p)
    for code, name in items:
        new_el = deepcopy(template_el)
        runs = new_el.findall(qn("a:r"))
        if len(runs) >= 2:
            t0 = runs[0].find(qn("a:t"))
            t1 = runs[1].find(qn("a:t"))
            if t0 is not None:
                t0.text = f"C{code} "
            if t1 is not None:
                t1.text = name
            for extra in runs[2:]:
                new_el.remove(extra)
        txbody.append(new_el)


def apply_criteria_split(prs, channel, values):
    """values : dict de build_participant_placeholders. Répartit tous les
    critères du canal en 2 groupes selon que leur note (ramenée sur 20)
    est >= (vert, gauche) ou < (rouge, droite) la note globale du canal
    pour ce participant (scope "vous" des deux côtés)."""
    slide_idx = CRITERIA_SLIDE_BY_CHANNEL.get(channel)
    if slide_idx is None or slide_idx >= len(prs.slides):
        return
    slide = prs.slides[slide_idx]
    green_box, red_box = _find_criteria_boxes(slide)
    if green_box is None or red_box is None:
        return

    channel_note_20 = _as_float(values.get(f"Total {channel} note vous"))
    if channel_note_20 is None:
        return

    green_items, red_items = [], []
    for code in CRITERIA_BY_CHANNEL[channel]:
        raw_note = _as_float(values.get(f"C{code} {channel} note vous"))
        if raw_note is None:
            continue
        weight = 2 if code == DOUBLED_CODE_BY_CHANNEL.get(channel) else 1
        note_20 = raw_note / (2 * weight) * 20
        name = CRITERION_SHORT_NAMES[channel][code]
        (green_items if note_20 >= channel_note_20 else red_items).append((code, name))

    _rebuild_criteria_list(green_box, green_items)
    _rebuild_criteria_list(red_box, red_items)


# ---------------------------------------------------------------- Point d'entrée

def apply_debrief_visuals(prs, participant, cache, values):
    """Point d'entrée unique, appelé après substitution des balises texte
    (voir app/restitutions/routes.py::create_restitution) : graphiques de
    comparaison/classement + classement texte + répartition des critères.
    Best-effort par élément : un élément dont les formes ne correspondent
    pas à la structure attendue est simplement ignoré, pour ne jamais faire
    échouer toute la génération à cause d'un seul graphique."""
    try:
        apply_comparison_charts(prs, values)
    except Exception:
        pass
    try:
        apply_ranking_charts(prs, participant, cache)
    except Exception:
        pass
    try:
        apply_classement_text(prs, participant, cache)
    except Exception:
        pass
    for channel in CHANNEL_ORDER:
        try:
            apply_criteria_split(prs, channel, values)
        except Exception:
            continue
