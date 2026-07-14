"""
Manipulation du fichier PowerPoint « Problématiques ».

La diapositive 2 (index 1) contient la liste des problématiques utilisée
comme inspiration pour la génération. Pour chaque scénario généré, on
ajoute une diapositive indiquant son numéro et l'URL où la réponse a été
trouvée (pas de capture d'écran, voir spécification).
"""

import io

from pptx import Presentation
from pptx.util import Inches, Pt

PROBLEMATIQUES_SLIDE_INDEX = 1


def read_problematiques_text(file_data):
    """Retourne le texte de toutes les formes de la diapositive 2."""
    prs = Presentation(io.BytesIO(file_data))
    if len(prs.slides) <= PROBLEMATIQUES_SLIDE_INDEX:
        return ""
    slide = prs.slides[PROBLEMATIQUES_SLIDE_INDEX]
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            texts.append(shape.text_frame.text.strip())
    return "\n".join(texts)


def append_scenario_slides(file_data, scenarios_with_numbers):
    """
    scenarios_with_numbers : liste de (numero_scenario, scenario_dict).
    Ajoute une diapositive par scénario avec son numéro et son url_source.
    Retourne les bytes du fichier mis à jour.
    """
    prs = Presentation(io.BytesIO(file_data))
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    for numero, scenario in scenarios_with_numbers:
        slide = prs.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        title_box.text_frame.text = f"Scénario {numero}"
        title_box.text_frame.paragraphs[0].font.size = Pt(28)
        title_box.text_frame.paragraphs[0].font.bold = True

        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(1.5))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.text = f"Source : {scenario.get('url_source', '')}"
        body_frame.paragraphs[0].font.size = Pt(16)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
