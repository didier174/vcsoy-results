"""
Génère un rapport d'étude (.pptx) à partir d'un modèle et d'un
dictionnaire de valeurs, en remplaçant les balises {{ Nom de la balise }}
dans tous les textes du modèle (titres, zones de texte, tableaux, formes
groupées).

Contrainte PowerPoint : le texte d'un même paragraphe visible peut être
réparti sur plusieurs "runs" XML (changement de police au milieu d'un mot,
correction automatique, etc.), donc une balise peut être coupée en deux
runs différents. On reconstitue le texte complet du paragraphe pour
repérer/remplacer les balises, puis on réécrit le résultat dans le
premier run (en vidant les suivants) pour ne pas perdre le formatage.

Coloration conditionnelle (vert/rouge, diapositives listées dans
COLOR_CODED_SLIDES) : confirmé, s'applique à toute cellule qui ne contient
QUE ("vous" ou "catégorie") comparée à son équivalent "tous" de la même
ligne — vert si supérieur, rouge si inférieur, aucune couleur en cas
d'égalité stricte ou de donnée manquante. On réutilise directement les
valeurs déjà formatées (pas de dictionnaire de valeurs brutes séparé) : la
diapositive n'a jamais qu'un seul type de grandeur par ligne (note, %, ou
durée), donc reparser le texte affiché suffit et évite de dupliquer tout
le calcul.
"""

import io
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

TAG_RE = re.compile(r"\{\{\s*([^{}]+?)\s*\}\}")
FULL_TAG_RE = re.compile(r"^\{\{\s*([^{}]+?)\s*\}\}$")

COLOR_CODED_SLIDES = {13, 15, 19, 21, 23, 27, 31}  # numéros 1-indexés
COLOR_GOOD = RGBColor(0x1E, 0x8E, 0x3E)
COLOR_BAD = RGBColor(0xC0, 0x1C, 0x28)

DURATION_RE = re.compile(r"^(?:(\d+)h)?(?:(\d+)min)?(?:(\d+)sec)?$")


def _comparable(value):
    """Reconvertit une valeur déjà formatée (note « 1,55 », pourcentage
    « 55% », durée « 2min54sec »/« 6h34min »/« 40sec », ou nombre brut) en
    float comparable. None si absente ("—") ou illisible."""
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    s = str(value).strip()
    if s in ("", "—", "-"):
        return None
    if s.endswith("%"):
        try:
            return float(s[:-1].replace(",", "."))
        except ValueError:
            return None
    m = DURATION_RE.match(s)
    if m and any(m.groups()):
        h, mi, se = (int(g) if g else 0 for g in m.groups())
        return h * 3600 + mi * 60 + se
    try:
        return float(s.replace(",", "."))
    except ValueError:
        return None


def _color_counterpart_tag(normalized_tag):
    """Le tag « ... vous »/« ... catégorie » -> le tag « ... tous »
    correspondant sur la même ligne, ou None si ce n'est pas un tag de ce
    type (ex. « ... tous »/« ... laureats » eux-mêmes, ou tag non numérique
    comme {{PARTICIPANT}})."""
    for suffix in (" vous", " categorie"):
        if normalized_tag.endswith(suffix):
            return normalized_tag[: -len(suffix)] + " tous"
    return None


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _iter_paragraphs(shape):
    if getattr(shape, "has_text_frame", False):
        yield from shape.text_frame.paragraphs
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                yield from cell.text_frame.paragraphs


def _paragraph_text(paragraph):
    return "".join(run.text for run in paragraph.runs)


def _normalize(tag):
    return tag.strip().lower()


def _substitute_paragraph(paragraph, lookup, color_enabled=False):
    full_text = _paragraph_text(paragraph)
    if "{{" not in full_text:
        return

    new_text = TAG_RE.sub(lambda m: str(lookup.get(_normalize(m.group(1)), m.group(0))), full_text)
    if new_text == full_text or not paragraph.runs:
        return

    paragraph.runs[0].text = new_text
    for run in paragraph.runs[1:]:
        run.text = ""

    if color_enabled:
        _apply_color(paragraph, lookup, full_text.strip())


def _apply_color(paragraph, lookup, original_text):
    """original_text : texte du paragraphe AVANT substitution. On ne
    colore que les cellules ne contenant RIEN d'autre qu'un seul tag
    {{...}} — le cas de toutes les cellules numériques des tableaux
    concernés — pour ne jamais toucher aux phrases de texte libre."""
    m = FULL_TAG_RE.match(original_text)
    if not m:
        return
    tag = _normalize(m.group(1))
    tous_tag = _color_counterpart_tag(tag)
    if tous_tag is None:
        return
    own = _comparable(lookup.get(tag))
    ref = _comparable(lookup.get(tous_tag))
    if own is None or ref is None or own == ref:
        return
    color = COLOR_GOOD if own > ref else COLOR_BAD
    for run in paragraph.runs:
        if run.text:
            run.font.color.rgb = color


def find_tags(template_bytes):
    """Retourne l'ensemble des noms de balises {{ ... }} présentes dans le modèle."""
    prs = Presentation(io.BytesIO(template_bytes))
    tags = set()
    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            for paragraph in _iter_paragraphs(shape):
                tags.update(TAG_RE.findall(_paragraph_text(paragraph)))
    return tags


def render_template(template_bytes, values):
    """
    template_bytes : contenu binaire du modèle .pptx.
    values : dict {nom de balise: valeur}. La correspondance balise/clé
    ignore la casse et les espaces superflus (ex. {{ PARTICIPANT }} et
    {{ participant }} pointent tous deux vers la clé "Participant").

    Retourne (bytes du .pptx généré, ensemble des balises inconnues
    rencontrées dans le modèle mais absentes de `values`).
    """
    prs = Presentation(io.BytesIO(template_bytes))
    lookup = {_normalize(key): value for key, value in values.items()}

    unknown = set()
    for slide_index, slide in enumerate(prs.slides, start=1):
        color_enabled = slide_index in COLOR_CODED_SLIDES
        for shape in _iter_shapes(slide.shapes):
            for paragraph in _iter_paragraphs(shape):
                text = _paragraph_text(paragraph)
                unknown.update(tag for tag in TAG_RE.findall(text) if _normalize(tag) not in lookup)
                _substitute_paragraph(paragraph, lookup, color_enabled=color_enabled)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue(), unknown
