"""
Met à jour les 2 graphiques natifs PowerPoint du modèle de rapport
d'étude, qui ne sont pas de simples balises {{ ... }} texte :

- diapositive 9 : jauge (barre empilée à 100%) note globale + note par
  canal, pour le participant, avec un rond de couleur positionné au bout
  de chaque barre et un texte "Canal non testé" pour un canal absent.
- diapositives 14/18/22/26/30 : mapping d'importance (nuage de points),
  un point par critère du canal — abscisse = taux de conformité de
  l'entreprise sur ce critère, ordonnée = importance de ce critère dans
  la note globale du canal (calculée sur l'ensemble de l'édition) — avec
  un calibrage des axes (min/max/croisement) adapté aux valeurs du
  participant, pour rester lisible comme dans le modèle.

Les balises {{ ... }} de generator.py ne s'appliquent qu'au texte : ces
graphiques nécessitent de modifier directement le cache de données XML
du graphique (repris tel quel par PowerPoint/LibreOffice à l'ouverture),
tout en conservant sa mise en forme (étiquettes, couleurs, position
manuelle des étiquettes, traits de rappel) déjà réglée dans le modèle.
"""

import copy
import math
import statistics
from itertools import permutations

from lxml import etree
from pptx.oxml.ns import qn

from app.models import TestResult
from app.results.presentation import CHANNEL_ORDER
from app.results.scoring import (
    build_compilation_rows,
    compute_criterion_stats,
)


# ------------------------------------------------------- Diapositive 9 : jauge

GAUGE_ROW_ORDER = ["global", "phone", "mail", "web", "rs", "chat"]
GAUGE_UNTESTED_TEXT_SHAPES = {"mail": "ZoneTexteMail", "rs": "ZoneTexteRes"}
GAUGE_OVAL_NAME_PREFIX = "Slide11_OvalNote"

# Position horizontale (EMU) d'un rond en fonction de la note sur 20 :
# régression linéaire sur les 6 ronds et leurs valeurs d'exemple dans le
# modèle d'origine (résidus < 50 000 EMU, soit < 1,5 mm — un très bon
# ajustement, qui confirme un positionnement proportionnel à la note sur
# toute la largeur de la jauge).
GAUGE_OVAL_X_SLOPE = 218753.91863099797
GAUGE_OVAL_X_INTERCEPT = 2824825.1388052916


def _gauge_oval_left(note_20):
    return round(GAUGE_OVAL_X_SLOPE * note_20 + GAUGE_OVAL_X_INTERCEPT)


def _set_hidden(shape, hidden):
    cnv_pr = shape._element.find(f".//{qn('p:cNvPr')}")
    if cnv_pr is None:
        return
    if hidden:
        cnv_pr.set("hidden", "1")
    elif "hidden" in cnv_pr.attrib:
        del cnv_pr.attrib["hidden"]


def _clone_untested_label(template_shape, target_top, new_id):
    new_el = copy.deepcopy(template_shape._element)
    cnv_pr = new_el.find(f".//{qn('p:cNvPr')}")
    cnv_pr.set("id", str(new_id))
    if "hidden" in cnv_pr.attrib:
        del cnv_pr.attrib["hidden"]
    off = new_el.find(f".//{qn('a:xfrm')}/{qn('a:off')}")
    off.set("y", str(target_top))
    return new_el


def apply_gauge_chart(prs, participant, edition_id, participant_tests=None):
    slide = prs.slides[8]  # diapositive 9
    chart_shape = next((s for s in slide.shapes if s.name == "Graph_Bar"), None)
    if chart_shape is None or not chart_shape.has_chart:
        return

    # La note consolidée et les notes par canal du participant ne dépendent
    # QUE de ses propres tests : pas besoin des données de toute l'édition.
    if participant_tests is None:
        participant_tests = TestResult.query.filter_by(edition_id=edition_id, participant_id=participant.id).all()
    own_row = build_compilation_rows([participant], participant_tests)[0]

    channel_flags = {
        "phone": participant.channel_phone, "mail": participant.channel_mail,
        "web": participant.channel_web, "rs": participant.channel_rs, "chat": participant.channel_chat,
    }

    notes = {}
    for key in GAUGE_ROW_ORDER:
        if key == "global":
            notes[key] = own_row["consolidated_score"] if own_row else None
        elif channel_flags.get(key):
            notes[key] = own_row["channels"][key]["note_20"] if own_row else None
        else:
            notes[key] = None

    series1 = [notes[key] if notes[key] is not None else 0.0 for key in GAUGE_ROW_ORDER]
    series2 = [max(0.0, 20.0 - v) for v in series1]

    all_ser = chart_shape.chart._chartSpace.findall(f".//{qn('c:ser')}")
    for ser_el, new_values in zip(all_ser, [series1, series2]):
        val_el = ser_el.find(qn("c:val"))
        numref = val_el.find(qn("c:numRef")) if val_el is not None else None
        numlit = val_el.find(qn("c:numLit")) if val_el is not None else None
        target = numlit if numlit is not None else (numref.find(qn("c:numCache")) if numref is not None else None)
        if target is not None:
            _set_numlit_points(target, {i: v for i, v in enumerate(new_values)})

    # Ronds de couleur : un par ligne (global/phone/mail/web/rs/chat), triés
    # par position verticale dans le modèle (2 d'entre eux partagent le même
    # nom "Slide11_OvalNoteReseau" par erreur de conception du modèle — on
    # se base donc sur l'ordre vertical, fiable, plutôt que sur le nom).
    ovals = sorted(
        (s for s in slide.shapes if s.name.startswith(GAUGE_OVAL_NAME_PREFIX)),
        key=lambda s: s.top,
    )
    for key, oval in zip(GAUGE_ROW_ORDER, ovals):
        note = notes[key]
        if note is None:
            _set_hidden(oval, True)
        else:
            _set_hidden(oval, False)
            oval.left = _gauge_oval_left(note)

    # "Canal non testé" à la place de la note : 2 zones de texte existent déjà
    # dans le modèle (Mail/RS, masquées par défaut) ; on les clone pour les 3
    # canaux qui n'en ont pas (Phone/Web/Chat), positionnées sur la même ligne
    # que leur rond (désormais masqué).
    label_template = next((s for s in slide.shapes if s.name == "ZoneTexteMail"), None)
    max_id = max((s.shape_id for s in slide.shapes), default=0)
    oval_by_row = dict(zip(GAUGE_ROW_ORDER, ovals))
    for channel in ("phone", "mail", "web", "rs", "chat"):
        tested = notes[channel] is not None
        shape_name = GAUGE_UNTESTED_TEXT_SHAPES.get(channel)
        if shape_name:
            shape = next((s for s in slide.shapes if s.name == shape_name), None)
            if shape is not None:
                _set_hidden(shape, tested)
        elif not tested and label_template is not None:
            oval = oval_by_row.get(channel)
            if oval is not None:
                target_top = oval.top + oval.height // 2 - label_template.height // 2
                max_id += 1
                new_el = _clone_untested_label(label_template, target_top, max_id)
                slide.shapes._spTree.append(new_el)


def _set_numlit_points(numlit_el, values_by_idx):
    for pt in list(numlit_el.findall(qn("c:pt"))):
        numlit_el.remove(pt)
    ptcount = numlit_el.find(qn("c:ptCount"))
    if ptcount is not None:
        ptcount.set("val", str(len(values_by_idx)))
    for idx in sorted(values_by_idx):
        pt = etree.SubElement(numlit_el, qn("c:pt"))
        pt.set("idx", str(idx))
        v = etree.SubElement(pt, qn("c:v"))
        v.text = repr(float(values_by_idx[idx]))


# --------------------------------------- Diapositives 14/18/22/26/30 : mapping

# idx (position du point dans le graphique, 0-based) -> code du critère
# "Code N" — déduit des étiquettes déjà présentes dans le modèle pour
# chaque canal (voir historique de la conversation). Le critère
# "Impression générale" est systématiquement absent du mapping sur les 5
# canaux (choix déjà fait dans le modèle d'origine, pas une balise à
# renseigner).
MAPPING_SLIDE_BY_CHANNEL = {"phone": 13, "mail": 17, "web": 21, "rs": 25, "chat": 29}  # 0-based
MAPPING_POINT_CODE = {
    "phone": {0: 1, 1: 2, 2: 3, 3: 4, 4: 5, 5: 6, 6: 7, 7: 8, 8: 9, 9: 10, 10: 11, 11: 12, 12: 13, 13: 15},
    "mail": {0: 1, 1: 2, 2: 3, 3: 4, 4: 5, 5: 6, 6: 7, 7: 8, 8: 9, 9: 11, 10: 10, 11: 13, 12: 14},
    "web": {0: 1, 1: 2, 2: 3, 3: 4, 4: 5, 5: 6, 6: 7, 7: 8, 8: 9, 9: 10, 10: 11, 11: 13},
    "rs": {0: 1, 1: 2, 2: 3, 3: 4, 4: 5, 5: 6, 6: 7, 7: 8, 8: 9, 9: 10, 10: 11, 11: 12, 12: 14, 13: 15},
    "chat": {0: 1, 1: 2, 2: 3, 3: 4, 4: 5, 5: 6, 6: 7, 7: 8, 8: 9, 9: 10, 10: 11, 11: 12, 12: 14},
}

# Marge (dans l'unité des données, 0-1 pour l'abscisse comme pour
# l'ordonnée) ajoutée de part et d'autre de la plage réelle des points pour
# calibrer les axes : déduite du modèle d'origine (mêmes échelles 0-1), où
# elle vaut ~0,3 des deux côtés. Le point de croisement des axes (souvent
# hors-centre dans le modèle) correspond, une fois vérifié, à la médiane des
# valeurs de l'AUTRE axe — ce qui place la ligne de croisement au milieu du
# nuage de points plutôt qu'à une valeur arbitraire comme 0.
MAPPING_AXIS_PADDING = 0.3


def _criterion_pct_vous(channel, code, vous_tests):
    """% de scores "Bon" (valeur brute = 2) du participant pour ce
    critère — même définition confirmée que pour les tableaux détaillés."""
    stats = compute_criterion_stats(channel, code, [t.raw_data for t in vous_tests])
    return (stats["pct"] or 0) / 100.0


def _scatter_axes(chartspace):
    """Retourne (axe_x, axe_y) : les 2 éléments <c:valAx> d'un graphique
    XY, identifiés via l'ordre de leurs <c:axId> dans <c:scatterChart>
    (confirmé, cet ordre est bien [abscisse, ordonnée]), pas leur ordre
    dans le document (moins fiable)."""
    scatter = chartspace.find(f".//{qn('c:scatterChart')}")
    if scatter is None:
        return None, None
    axid_els = scatter.findall(qn("c:axId"))
    if len(axid_els) < 2:
        return None, None
    x_axid, y_axid = axid_els[0].get("val"), axid_els[1].get("val")
    by_id = {}
    for ax in chartspace.findall(f".//{qn('c:valAx')}"):
        axid_el = ax.find(qn("c:axId"))
        if axid_el is not None:
            by_id[axid_el.get("val")] = ax
    return by_id.get(x_axid), by_id.get(y_axid)


def _set_axis_range(axis_el, data_min, data_max, crosses_at):
    if axis_el is None:
        return
    # c:min/c:max sont dans c:scaling (petit-enfant de c:valAx, pas enfant
    # direct) : recherche en profondeur nécessaire, contrairement à
    # c:crossesAt qui est bien un enfant direct de c:valAx.
    min_el = axis_el.find(f".//{qn('c:min')}")
    max_el = axis_el.find(f".//{qn('c:max')}")
    crosses_el = axis_el.find(qn("c:crossesAt"))
    if min_el is not None:
        min_el.set("val", repr(data_min - MAPPING_AXIS_PADDING))
    if max_el is not None:
        max_el.set("val", repr(data_max + MAPPING_AXIS_PADDING))
    if crosses_el is not None:
        crosses_el.set("val", repr(crosses_at))


# Modèle géométrique (unités normalisées 0-1) d'une zone de texte : une
# boîte centrée sur son point d'ancrage (position du point + décalage),
# largeur estimée à partir du nombre de caractères réel de l'étiquette
# (6 à 44 caractères selon les critères — une largeur fixe sous-estimait
# largement le besoin des plus longues), hauteur fixe (texte sur 1 ligne).
LABEL_HEIGHT = 0.06
LABEL_CHAR_WIDTH = 0.006
LABEL_MARKER_WIDTH = 0.04
# Marge visible entre 2 boîtes de texte (en plus de leur non-chevauchement
# strict) et pas minimal forcé entre étiquettes consécutives d'un même
# groupe (pour garantir un ordre vertical identique à celui des points,
# condition nécessaire pour que les traits de rappel ne se croisent pas).
OVERLAP_BUFFER = 0.015
ORDER_MARGIN = 0.003


def _label_text(dlbl):
    return "".join(t.text or "" for t in dlbl.iter(qn("a:t")))


def _label_width(dlbl):
    return LABEL_MARKER_WIDTH + LABEL_CHAR_WIDTH * len(_label_text(dlbl))


def _dlbl_layout_offset(dlbl):
    layout = dlbl.find(qn("c:layout"))
    ml = layout.find(qn("c:manualLayout")) if layout is not None else None
    x_el = ml.find(qn("c:x")) if ml is not None else None
    y_el = ml.find(qn("c:y")) if ml is not None else None
    dx = float(x_el.get("val")) if x_el is not None else 0.0
    dy = float(y_el.get("val")) if y_el is not None else 0.0
    return dx, dy


def _set_dlbl_layout_offset(dlbl, dx, dy):
    layout = dlbl.find(qn("c:layout"))
    if layout is None:
        # <c:layout> doit être placé juste après <c:idx> (ordre imposé par
        # le schéma), avant <c:tx>/<c:spPr>/... — jamais en fin d'élément.
        layout = etree.Element(qn("c:layout"))
        dlbl.find(qn("c:idx")).addnext(layout)
    ml = layout.find(qn("c:manualLayout"))
    if ml is None:
        ml = etree.SubElement(layout, qn("c:manualLayout"))
    x_el = ml.find(qn("c:x"))
    if x_el is None:
        x_el = etree.SubElement(ml, qn("c:x"))
    x_el.set("val", repr(dx))
    y_el = ml.find(qn("c:y"))
    if y_el is None:
        y_el = etree.SubElement(ml, qn("c:y"))
    y_el.set("val", repr(dy))


# Marge (mêmes unités normalisées 0-1) gardée entre une étiquette et la
# ligne de croisement des axes, pour qu'elle reste nettement à l'intérieur
# de son quadrant plutôt que juste effleurer la limite. Sert aussi de garde
# minimale entre 2 étiquettes de part et d'autre de la ligne (2x cette
# marge), le cas le plus difficile puisqu'on ne peut pas les rapprocher
# davantage sans changer l'une d'elles de quadrant.
QUADRANT_MARGIN = 0.04

MAX_UNCROSS_PASSES = 6


def _segments_intersect(p1, p2, p3, p4):
    def ccw(a, b, c):
        return (c[1] - a[1]) * (b[0] - a[0]) > (b[1] - a[1]) * (c[0] - a[0])
    return ccw(p1, p3, p4) != ccw(p2, p3, p4) and ccw(p1, p2, p3) != ccw(p1, p2, p4)


def _count_crossings(idxs, point_of, ax, ay):
    total = 0
    for i in range(len(idxs)):
        a = idxs[i]
        for b in idxs[i + 1:]:
            if _segments_intersect(point_of[a], (ax[a], ay[a]), point_of[b], (ax[b], ay[b])):
                total += 1
    return total


def _boxes_overlap(ax_a, ay_a, width_a, ax_b, ay_b, width_b):
    # width_a/width_b et LABEL_HEIGHT sont des demi-étendues (voir
    # clamp_box) : 2 boîtes ne se chevauchent pas si distantes d'au moins
    # la SOMME de leurs demi-étendues respectives, pas leur moyenne.
    return (
        abs(ax_a - ax_b) < width_a + width_b + OVERLAP_BUFFER
        and abs(ay_a - ay_b) < 2 * LABEL_HEIGHT + OVERLAP_BUFFER
    )


def _has_any_overlap(idxs, ax, final_y, width_of):
    for i in range(len(idxs)):
        a = idxs[i]
        for b in idxs[i + 1:]:
            if _boxes_overlap(ax[a], final_y[a], width_of[a], ax[b], final_y[b], width_of[b]):
                return True
    return False


def _quadrant_ok(ax, ay, width, nx, ny, cx, cy):
    if nx >= cx:
        if ax - width < cx + QUADRANT_MARGIN - 1e-9:
            return False
    elif ax + width > cx - QUADRANT_MARGIN + 1e-9:
        return False
    if ny >= cy:
        if ay - LABEL_HEIGHT < cy + QUADRANT_MARGIN - 1e-9:
            return False
    elif ay + LABEL_HEIGHT > cy - QUADRANT_MARGIN + 1e-9:
        return False
    return True


# Taille de groupe (par quadrant) en dessous de laquelle on essaie TOUTES
# les permutations possibles des emplacements (garantit la meilleure
# solution possible) ; au-delà, une recherche gloutonne par paires (plus
# rapide mais pas nécessairement optimale) prend le relais — 8! = 40 320
# reste instantané, mais un groupe plus grand rendrait le balayage complet
# trop coûteux pour un simple raffinement visuel.
EXHAUSTIVE_UNCROSS_LIMIT = 8


def _resolve_leader_crossings(points, cx, cy, ax, final_y):
    """Les traits de rappel (point -> étiquette) ne doivent jamais se
    croiser entre eux (confirmé indispensable) : l'ordre vertical imposé
    plus haut l'évite déjà pour la plupart des cas, mais 2 points très
    proches l'un de l'autre peuvent quand même produire des traits croisés
    selon leurs décalages horizontaux respectifs (hérités du modèle).

    On ne peut réattribuer les emplacements (ax, ay) de 2 étiquettes qu'à
    l'intérieur d'un même quadrant (même position par rapport aux 2 lignes
    de croisement). Comme les étiquettes réattribuées n'ont pas forcément
    la même largeur, un emplacement valide pour l'une ne l'est pas
    forcément pour l'autre (règles 1-2, prioritaires) — on revérifie donc
    le quadrant après réattribution, en plus de l'absence de chevauchement
    (règle 3), et on ne retient que l'arrangement qui minimise le nombre de
    croisements du quadrant parmi tous ceux qui respectent ces 2 règles."""
    quadrants = {}
    for idx, nx, ny, _, _, above, _ in points:
        quadrants.setdefault((above, nx >= cx), []).append(idx)
    point_of = {p[0]: (p[1], p[2]) for p in points}
    width_of = {p[0]: p[6] for p in points}

    for idxs in quadrants.values():
        n = len(idxs)
        if n < 2:
            continue
        if n <= EXHAUSTIVE_UNCROSS_LIMIT:
            _uncross_exhaustive(idxs, point_of, width_of, cx, cy, ax, final_y)
        else:
            _uncross_greedy(idxs, point_of, width_of, cx, cy, ax, final_y)


def _uncross_exhaustive(idxs, point_of, width_of, cx, cy, ax, final_y):
    n = len(idxs)
    slots = [(ax[idx], final_y[idx]) for idx in idxs]
    identity = tuple(range(n))
    best_perm = identity
    best_count = _count_crossings(idxs, point_of, ax, final_y)
    if best_count == 0:
        return
    for perm in permutations(range(n)):
        if perm == identity:
            continue
        cand_ax = {idxs[i]: slots[perm[i]][0] for i in range(n)}
        cand_ay = {idxs[i]: slots[perm[i]][1] for i in range(n)}
        if not all(
            _quadrant_ok(cand_ax[idx], cand_ay[idx], width_of[idx], *point_of[idx], cx, cy)
            for idx in idxs
        ):
            continue
        if _has_any_overlap(idxs, cand_ax, cand_ay, width_of):
            continue
        count = _count_crossings(idxs, point_of, cand_ax, cand_ay)
        if count < best_count:
            best_count, best_perm = count, perm
            if best_count == 0:
                break
    if best_perm != identity:
        for i, idx in enumerate(idxs):
            ax[idx] = slots[best_perm[i]][0]
            final_y[idx] = slots[best_perm[i]][1]


def _uncross_greedy(idxs, point_of, width_of, cx, cy, ax, final_y):
    for _ in range(MAX_UNCROSS_PASSES):
        base = _count_crossings(idxs, point_of, ax, final_y)
        if base == 0:
            break
        improved = False
        for i in range(len(idxs)):
            a = idxs[i]
            for b in idxs[i + 1:]:
                ax[a], ax[b] = ax[b], ax[a]
                final_y[a], final_y[b] = final_y[b], final_y[a]
                valid = (
                    _quadrant_ok(ax[a], final_y[a], width_of[a], *point_of[a], cx, cy)
                    and _quadrant_ok(ax[b], final_y[b], width_of[b], *point_of[b], cx, cy)
                    and not _has_any_overlap(idxs, ax, final_y, width_of)
                )
                new_total = _count_crossings(idxs, point_of, ax, final_y)
                if valid and new_total < base:
                    base = new_total
                    improved = True
                else:
                    ax[a], ax[b] = ax[b], ax[a]
                    final_y[a], final_y[b] = final_y[b], final_y[a]
        if not improved:
            break


def _spread_overlapping_labels(ser_el, x_by_idx, y_by_idx, x_min, x_max, y_min, y_max, x_crosses, y_crosses):
    """Repositionne les étiquettes (boîtes rectangulaires centrées sur leur
    point d'ancrage, voir LABEL_HEIGHT/_label_width) pour satisfaire, dans
    l'ordre de priorité suivant — confirmé indispensable pour la lecture du
    mapping :
    1. la boîte ENTIÈRE (pas seulement son point de départ) reste du même
       côté des lignes de croisement des axes (x_crosses/y_crosses) que son
       point de données — même quadrant, de bout en bout ;
    2. aucune boîte ne chevauche une autre (test rectangle strict, avec une
       marge visible) ;
    3. l'ordre vertical des étiquettes reste identique à celui de leurs
       points (au sein d'un même demi-plan haut/bas) : les traits de rappel
       ne peuvent alors pas se croiser entre eux, puisque 2 segments reliant
       des paires ordonnées de la même façon aux deux extrémités ne se
       croisent jamais.

    Le décalage horizontal déjà réglé dans le modèle est d'abord borné pour
    que la boîte entière reste du bon côté de la ligne verticale.
    Verticalement, les étiquettes "au-dessus" et "en dessous" de la ligne
    horizontale sont balayées séparément (chaque groupe comparé uniquement
    à lui-même, dans l'ordre de leur position d'origine), en poussant
    toujours À L'ÉCART de la ligne plutôt que vers elle : le quadrant est
    donc respecté par construction, jamais par un plafond appliqué après
    coup qui laisserait des étiquettes se coincer près de la ligne.

    Le trait de rappel (déjà activé dans le modèle, voir showLeaderLines)
    suit alors automatiquement l'étiquette jusqu'à sa nouvelle position —
    PowerPoint le calcule à l'affichage, pas besoin de le dessiner ici.

    Piège : c:manualLayout x/y est exprimé dans le repère de mise en page
    du GRAPHIQUE (haut de la zone = 0, vers le bas = positif), l'inverse du
    repère des DONNÉES qu'on utilise ici pour nx/ny (valeur qui augmente =
    vers le haut). On travaille donc en interne avec vy = -dy (repère
    "visuel", cohérent avec ny), et on réinverse le signe au moment
    d'écrire la valeur dans le XML."""
    dlbls_el = ser_el.find(qn("c:dLbls"))
    if dlbls_el is None:
        return
    dlbl_by_idx = {}
    for dlbl in dlbls_el.findall(qn("c:dLbl")):
        idx_el = dlbl.find(qn("c:idx"))
        if idx_el is not None:
            dlbl_by_idx[int(idx_el.get("val"))] = dlbl

    def norm(v, lo, hi):
        return (v - lo) / (hi - lo) if hi > lo else 0.5

    cx = norm(x_crosses, x_min, x_max)
    cy = norm(y_crosses, y_min, y_max)

    def clamp_box(value, crossing, is_positive_side, extent):
        # extent : la largeur/hauteur ENTIÈRE de la boîte, pas sa moitié —
        # l'alignement du texte par rapport à son point d'ancrage
        # (c:dLblPos) n'étant pas connu avec certitude, on se protège du
        # pire cas où la boîte s'étend entièrement d'un seul côté de
        # l'ancre plutôt que de part et d'autre.
        if is_positive_side:
            return max(value, crossing + QUADRANT_MARGIN + extent)
        return min(value, crossing - QUADRANT_MARGIN - extent)

    points = []
    for idx in x_by_idx:
        if idx not in dlbl_by_idx:
            continue
        nx = norm(x_by_idx[idx], x_min, x_max)
        ny = norm(y_by_idx[idx], y_min, y_max)
        dx0, dy0 = _dlbl_layout_offset(dlbl_by_idx[idx])
        vy0 = -dy0  # repère "visuel" (haut = grand), voir note ci-dessus
        width = _label_width(dlbl_by_idx[idx])
        ax = clamp_box(nx + dx0, cx, nx >= cx, width)
        points.append((idx, nx, ny, ax, vy0, ny >= cy, width))
    if not points:
        return

    # Balayage séparé par quadrant (haut/bas x droite/gauche de cx/cy) :
    # chaque groupe n'est comparé qu'à lui-même, dans l'ordre de la
    # position d'origine (ny) de ses points. Verticalement, on pousse
    # toujours À L'ÉCART de la ligne (jamais vers elle), ce qui garantit le
    # respect du quadrant par construction. Horizontalement, on force en
    # plus l'ordre des abscisses à suivre celui des ordonnées : 2 segments
    # reliant des paires dont l'ordre est identique sur les DEUX axes ne
    # peuvent pas se croiser, alors que ne préserver que l'ordre vertical
    # ne suffit pas (un décalage horizontal hérité du modèle peut à lui
    # seul inverser 2 traits de rappel).
    final_y = {}
    ax_by_idx = {}
    for above in (True, False):
        for right in (True, False):
            group = [p for p in points if p[5] == above and (p[1] >= cx) == right]
            if not group:
                continue
            dir_y = 1 if above else -1
            dir_x = 1 if right else -1
            # Trié par angle depuis le point de croisement (cx, cy), pas
            # seulement par ny : ordonner les étiquettes selon un seul axe
            # (comme avant) ne suffit pas à empêcher un croisement si les
            # abscisses des points ne suivent pas elles-mêmes cet ordre —
            # l'angle capture les 2 axes à la fois dans un seul critère.
            group.sort(key=lambda p: math.atan2(dir_y * (p[2] - cy), dir_x * (p[1] - cx)))
            placed = []
            for idx, nx, ny, ax, vy0, _, width in group:
                if placed:
                    # Ordre horizontal identique à celui des points (avant
                    # le calcul de l'anti-chevauchement, qui doit porter sur
                    # la position ax définitive) : sans quoi l'ordre
                    # vertical seul ne suffit pas à empêcher un croisement
                    # (voir docstring, point 3).
                    prev_ax = placed[-1][0]
                    ax = max(ax, prev_ax) if dir_x > 0 else min(ax, prev_ax)
                y = clamp_box(ny + vy0, cy, above, LABEL_HEIGHT)
                for other_x, other_y, other_width in placed:
                    # width/other_width et LABEL_HEIGHT sont des demi-étendues
                    # (voir clamp_box) : la marge de non-chevauchement est
                    # donc leur SOMME, pas leur moyenne.
                    if abs(ax - other_x) < width + other_width + OVERLAP_BUFFER:
                        min_gap = 2 * LABEL_HEIGHT + OVERLAP_BUFFER
                        y = max(y, other_y + min_gap) if dir_y > 0 else min(y, other_y - min_gap)
                if placed:
                    # Ordre vertical identique à celui des points (voir
                    # docstring, point 3).
                    prev_y = placed[-1][1]
                    y = max(y, prev_y + ORDER_MARGIN) if dir_y > 0 else min(y, prev_y - ORDER_MARGIN)
                placed.append((ax, y, width))
                final_y[idx] = y
                ax_by_idx[idx] = ax

    _resolve_leader_crossings(points, cx, cy, ax_by_idx, final_y)

    for idx, nx, ny, ax, vy0, above, width in points:
        dlbl = dlbl_by_idx[idx]
        final_ax = ax_by_idx[idx]
        # -(...) : reconversion du repère "visuel" vers celui, inversé, du
        # c:manualLayout XML (voir note en tête de fonction).
        _set_dlbl_layout_offset(dlbl, final_ax - nx, -(final_y[idx] - ny))


def apply_importance_mappings(prs, participant, edition_id, cache, participant_tests=None):
    """cache["importance"][channel][str(code)] : coefficient de Pearson au
    carré, précalculé une fois pour toute l'édition (voir report_cache.py)
    — évite de recharger/rescanner tous les tests de l'édition à chaque
    génération de rapport."""
    if participant_tests is None:
        participant_tests = TestResult.query.filter_by(edition_id=edition_id, participant_id=participant.id).all()
    vous_tests_by_channel = {
        channel: [t for t in participant_tests if t.channel == channel]
        for channel in CHANNEL_ORDER
    }

    for channel in CHANNEL_ORDER:
        slide_idx = MAPPING_SLIDE_BY_CHANNEL[channel]
        slide = prs.slides[slide_idx]
        chart_shape = next(
            (s for s in slide.shapes if getattr(s, "has_chart", False) and s.chart.chart_type == -4169),
            None,
        )
        if chart_shape is None:
            continue

        point_codes = MAPPING_POINT_CODE[channel]

        # Importance = coefficient de Pearson au carré (toujours positif,
        # reflète l'intensité de la liaison indépendamment du sens), normalisé
        # pour que la somme sur tous les critères du canal fasse 1 — même
        # échelle que les valeurs d'exemple du modèle d'origine.
        raw_importance = cache["importance"][channel]
        total = sum(raw_importance[str(code)] for code in point_codes.values())

        x_by_idx, y_by_idx = {}, {}
        vous_tests = vous_tests_by_channel[channel]
        for idx, code in point_codes.items():
            x_by_idx[idx] = _criterion_pct_vous(channel, code, vous_tests)
            y_by_idx[idx] = (raw_importance[str(code)] / total) if total else 0.0

        chart_xml = chart_shape.chart._chartSpace
        ser_el = chart_xml.find(f".//{qn('c:ser')}")
        if ser_el is None:
            continue
        xval_el = ser_el.find(qn("c:xVal"))
        yval_el = ser_el.find(qn("c:yVal"))
        if xval_el is not None:
            numlit = xval_el.find(qn("c:numLit"))
            if numlit is not None:
                _set_numlit_points(numlit, x_by_idx)
        if yval_el is not None:
            numlit = yval_el.find(qn("c:numLit"))
            if numlit is not None:
                _set_numlit_points(numlit, y_by_idx)

        # Calibrage des axes sur la plage réelle des valeurs du participant :
        # sans ça, les axes restent calés sur les valeurs d'exemple du
        # modèle, ce qui peut resserrer/décentrer le nuage de points et faire
        # se chevaucher les étiquettes (déjà positionnées manuellement dans
        # le modèle pour SA plage de valeurs d'origine).
        if x_by_idx and y_by_idx:
            x_axis, y_axis = _scatter_axes(chart_xml)
            xs, ys = list(x_by_idx.values()), list(y_by_idx.values())
            x_min, x_max = min(xs) - MAPPING_AXIS_PADDING, max(xs) + MAPPING_AXIS_PADDING
            y_min, y_max = min(ys) - MAPPING_AXIS_PADDING, max(ys) + MAPPING_AXIS_PADDING
            x_crosses, y_crosses = statistics.median(xs), statistics.median(ys)
            _set_axis_range(x_axis, min(xs), max(xs), y_crosses)
            _set_axis_range(y_axis, min(ys), max(ys), x_crosses)
            _spread_overlapping_labels(
                ser_el, x_by_idx, y_by_idx, x_min, x_max, y_min, y_max, x_crosses, y_crosses
            )


def apply_report_visuals(prs, participant, edition_id, cache, participant_tests=None):
    """Point d'entrée unique : applique la jauge (diapo 9) et les 5
    mappings d'importance (diapos 14/18/22/26/30) sur une Presentation déjà
    ouverte (après substitution des balises texte).

    cache : agrégats d'édition précalculés (voir report_cache.py).
    participant_tests : les tests du SEUL participant courant, déjà
    calculés par l'appelant (voir reports/routes.py) pour éviter une
    requête redondante."""
    apply_gauge_chart(prs, participant, edition_id, participant_tests=participant_tests)
    apply_importance_mappings(prs, participant, edition_id, cache, participant_tests=participant_tests)
